# Importing required libraries for logging and others
import logging
import io
import sys
import openai
import streamlit as st
import PyPDF2
import docx
from typing import Union
from pptx import Presentation
from pydub import AudioSegment
import requests
import time
import tiktoken

# Streamlit title
st.title("Josh's AI Assistant")

print("This is a test")
# Configure API key for OpenAI
openai.api_key = st.secrets["openai"]["api_key"]

# Configure logging to print to console
logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s]: %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)]
)


# Function to count the number of tokens in a string
def num_tokens_from_string(string: str, encoding_name: str = "cl100k_base") -> int:
    encoding = tiktoken.get_encoding(encoding_name)
    num_tokens = len(encoding.encode(string))
    return num_tokens

# Constants for file types
PDF = "application/pdf"
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPT = "application/vnd.ms-powerpoint"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
TXT = "text/plain"

# Function to handle PDF files
def handle_pdf(file):
    logging.debug("Handling PDF file.")
    pdf_reader = PyPDF2.PdfReader(file)
    return [pdf_reader.pages[page_num].extract_text() for page_num in range(len(pdf_reader.pages))]

# Function to handle DOCX files
def handle_docx(file):
    logging.debug("Handling DOCX file.")
    doc = docx.Document(io.BytesIO(file.read()))
    return ["\n".join([para.text for para in doc.paragraphs])]

# Function to handle PPTX files
def handle_pptx(file):
    logging.debug("Handling PPTX file.")
    prs = Presentation(io.BytesIO(file.read()))
    return [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]

# Function to handle TXT files
def handle_txt(file):
    logging.debug("Handling TXT file.")
    return [file.read().decode("utf-8")]

# Function to handle the uploaded file and extract text
def handle_uploaded_file(uploaded_file: Union[io.IOBase, object]) -> str:
    logging.debug("Handling uploaded file.")
    text_list = []
    file_type = uploaded_file.type
    if file_type == PDF:
        text_list.extend(handle_pdf(uploaded_file))
    elif file_type == DOCX:
        text_list.extend(handle_docx(uploaded_file))
    elif file_type in [PPT, PPTX]:
        text_list.extend(handle_pptx(uploaded_file))
    elif file_type == TXT:
        text_list.extend(handle_txt(uploaded_file))
    else:
        logging.error("Unsupported file type.")
        return "Unsupported file type"
    return f"Document: {''.join(text_list)}"

# Function to handle audio data
def handle_audio_data(uploaded_file):
    logging.debug("Handling audio data.")
    audio_data = None
    if uploaded_file is not None and uploaded_file.type in ["audio/mp3", "audio/wav", "audio/m4a"]:
        audio_data = uploaded_file.read()
    if audio_data is not None:
        file_format = uploaded_file.type.split('/')[-1]
        audio_segment = AudioSegment.from_file(io.BytesIO(audio_data), format=file_format)
        audio_segment = audio_segment.set_frame_rate(16000)
        audio_file = io.BytesIO()
        audio_segment.export(audio_file, format="wav")
        audio_file.seek(0)
        audio_file.name = "audio.wav"
        return audio_file

def transcribe_audio(audio_file):
    logging.debug("Transcribing audio.")
    assembly_key = st.secrets["assemblyai"]["api_key"]
    headers = {"authorization": assembly_key, "content-type": "application/json"}
    upload_endpoint = "https://api.assemblyai.com/v2/upload"
    transcription_endpoint = "https://api.assemblyai.com/v2/transcript"

    audio_file.seek(0)  # Reset the file pointer to the beginning
    response = requests.post(upload_endpoint, headers=headers, data=audio_file.read())
    audio_url = response.json()["upload_url"]

    payload = {"audio_url": audio_url}
    response = requests.post(transcription_endpoint, headers=headers, json=payload)
    job_id = response.json()["id"]

    while True:
        response = requests.get(f"{transcription_endpoint}/{job_id}", headers=headers)
        status = response.json()["status"]
        if status == "completed":
            transcript = response.json()["text"]
            break
        elif status == "error":
            transcript = "Error occurred during transcription."
            break
        time.sleep(5)

    return transcript

# Function to batch messages for OpenAI API
def batch_messages(messages, max_tokens=16385):
    logging.debug("Starting to batch messages.")
    batches = []
    current_batch = []
    current_token_count = 0
    for message in messages:
        message_tokens = num_tokens_from_string(message["content"])
        logging.debug(f"Message tokens: {message_tokens}, Current batch tokens: {current_token_count}")
        if current_token_count + message_tokens > max_tokens:
            logging.debug(f"Creating new batch due to token limit. Current batch tokens: {current_token_count}")
            batches.append(current_batch)
            current_batch = []
            current_token_count = 0
        current_batch.append(message)
        current_token_count += message_tokens
    if current_batch:
        logging.debug(f"Finalizing last batch with tokens: {current_token_count}")
        batches.append(current_batch)
    return batches
    
# Function to handle chat interaction
def handle_chat(prompt, context_document):
    logging.debug("Starting chat interaction.")
    if "openai_model" not in st.session_state:
        st.session_state["openai_model"] = "gpt-3.5-turbo-16k"
    if "messages" not in st.session_state:
        st.session_state.messages = []
    total_tokens = sum(num_tokens_from_string(msg["content"]) for msg in st.session_state.messages)
    logging.debug(f"Total tokens in all messages: {total_tokens}")

    # Display only user and assistant messages in chat interface
    for message in st.session_state.messages:
        if message["role"] in ["user", "assistant"]:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    if prompt:
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # In handle_chat function
        temp_messages = st.session_state.messages + [{"role": "system", "content": final_context}]
        logging.debug(f"Temp messages: {temp_messages}")

        
        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            full_response = ""
            message_batches = batch_messages(temp_messages)
            for batch in message_batches:
                for response in openai.ChatCompletion.create(
                    model=st.session_state["openai_model"],
                    messages=batch,
                    stream=True,
                ):
                    full_response += response.choices[0].delta.get("content", "")
                    message_placeholder.markdown(full_response + "▌", unsafe_allow_html=True)
            message_placeholder.markdown(full_response, unsafe_allow_html=True)
        
        st.session_state.messages.append({"role": "assistant", "content": full_response})


# Streamlit session state initialization
if "messages" not in st.session_state:
    st.session_state.messages = []

# File upload interface
uploaded_file = st.file_uploader("Upload your document or audio file here", type=["pdf", "docx", "ppt", "pptx", "txt", "mp3", "wav", "m4a"])

# Handle uploaded file and audio data
context_document = ""
audio_transcript = ""
if uploaded_file is not None:
    # Handling document files
    if uploaded_file.type in ["pdf", "docx", "ppt", "pptx", "txt"]:
        context_document = handle_uploaded_file(uploaded_file)
    # Handling audio files
    elif uploaded_file.type in ["audio/mp3", "audio/wav", "audio/m4a"]:
        audio_file = handle_audio_data(uploaded_file)
        if audio_file is not None:
            audio_transcript = transcribe_audio(audio_file)
            st.write(f"Audio Transcript: {audio_transcript}")

# Merge document context and audio transcript
final_context = context_document + "\n" + audio_transcript

logging.debug(f"Final context: {final_context}")
st.write(f"Debug: Final context is {final_context}")

# Chat interface
prompt = st.chat_input("What is this document about?")
handle_chat(prompt, final_context)

# Show logging output
st.write("Logging Output")
#st.text_area("", log_stream.getvalue())

